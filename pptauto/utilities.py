from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import pptx
import os
import unittest

def get_balloon_numbers(slide):
    """ Returns a list with the numbers inside the balloons in the slide """

    numbers = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text != '':
            numbers.append(shape.text) # ellipse
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shape in shape.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text != '':
                    numbers.append(shape.text) # ellipse inside a group
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shape in shape.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                            numbers.append(shape.text) # text box inside group inside group
                            
    numbers = [int(number.strip()) for number in numbers]
    numbers = set(numbers) # removes duplicates and sorts numbers
    numbers_str = [str(number) for number in numbers]
    return numbers_str

def write_with_format(cell, text, font_name = 'Arial', font_size = 10):
    """ Writes the string text in the cell using font_name with font_size"""
    text_frame = cell.text_frame
    text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    
def write_numbers_in_table(table, numbers):
    """ Writes (with format) the numbers in the first column of the table"""
    row = 2
    for number in numbers:
        cell = table.cell(row,0)
        write_with_format(cell, number)
        row += 1
        
def get_table(slide):
    """Returns the BOM table of the slide"""
    for shape in slide.shapes:
        if not shape.has_table or shape.table.cell(0,0).text != 'LISTA DE PARTES':
            continue
        return shape.table
    
def number_paragraphs(paragraphs, initial_number = 1):
    """Numbers the paragraphs given"""
    for paragraph in paragraphs:
        aux_list = paragraph.text.split('.')
        paragraph.text = '.'.join(str(initial_number), aux_list[1])
        initial_number += 1

def get_steps_paragraphs(slides):
    """Returns a list with the paragraphs that describe a step in the slides"""
    results = []
    for slide in slides:
        shape = slide.placeholders[12]
        text_frame = shape.text_frame
        paragraphs_read = text_frame.paragraphs
        first_par = paragraphs_read[0]

        if first_par.text.count('NOTA') +  first_par.text.count('ATENCION') +  first_par.text.count('-') != 0:
            results.append(first_par)

        for paragraph in paragraphs_read:
            if paragraph.text.count('NOTA') +  paragraph.text.count('ATENCION') +  paragraph.text.count('-'):
                continue
            results.append(paragraph)
    
    return results

def get_task_dict(prs):
    """Returns a dictionary with the task numbers as keys and a list with the slides corresponding to the task number for each key"""

    slides_tasks = prs.slide_layouts[4].used_by_slides
    task_dict = {}

    for slide in slides_tasks:
        shape = slide.placeholders[10]
        task_id = shape.text.strip()
        task_dict.setdefault(task_id,[]).append(slide)

    return task_dict

def write_task_number_text (prs, task_number, task_text, start, end):
    """Writes the task number and text from slide number start to slide number end """
    for i in range(start, end+1):
        working_slide = prs.slides[i]
        working_slide.placeholders[10].text = task_number.upper()
        working_slide.placeholders[11].text = task_text.upper()

def get_slides_task(prs):
    return prs.slide_layouts[4].used_by_slides

def fill_BOM (prs):
    slides_tasks = get_slides_task(prs)
    
    # Filling the BOM table in each slide:

    for slide in slides_tasks:
        numbers = get_balloon_numbers(slide)
        table = get_table(slide)
        write_numbers_in_table(table,numbers) 

def file_type(string):
    """Returns True if the string passed is a path to a file. It is an auxiliary function to check that input and output flags are paths to a file"""
    if os.path.isfile(string):
        return(string)
    else:
        raise FileNotFoundError(string)
    
class ModuleTestCase(unittest.TestCase):
    
    def setUp(self):
        self.prs = pptx.Presentation('data/IT-prueba.pptx')

    def test_get_slides_task(self):
        slides_task = get_slides_task(self.prs)
        self.assertEqual(len(slides_task), 14)

    def test_get_balloon_numbers(self):
        slide = self.prs.slides[5]
        self.assertEqual(get_balloon_numbers(slide), ['1','5','6','7'])

    def test_get_table(self):
        slide = self.prs.slides[5]
        self.assertTrue(isinstance(get_table(slide), pptx.table.Table))

if __name__ == '__main__':
    unittest.main()
    # python -m pptauto.utilities to run the tests